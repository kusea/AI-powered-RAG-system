#Dỉrectly work with SQLAlchemy to store and query vector
from sqlalchemy.orm import Session
from sqlalchemy import or_
from app.models import Document, ChunkDocument, User, DocumentShare, Notification, DocumentInsight
from app.schemas.document import ChunkEmbeddingCreate, DocumentShareCreate
from app.core.notification_mannager import notification_manager
from app.services import rag_service
from app.core.database import SessionLocal
from sentence_transformers import SentenceTransformer

import os 
import uuid
import pandas as pd
import pytesseract
import pdfplumber
import io

from PIL import Image
from datetime import datetime, timezone
from shutil import copyfileobj
from fastapi import UploadFile, HTTPException
from docx import Document as DocxReader 
from pptx import Presentation
from langchain_text_splitters import MarkdownHeaderTextSplitter, RecursiveCharacterTextSplitter
from googleapiclient.discovery import build
from google.oauth2.credentials import Credentials
from googleapiclient.http import MediaIoBaseDownload


# Encode content (or title) text to vector embedding
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

pytesseract.pytesseract.tesseract_cmd = r'D:\Tesseract-OCR\tesseract.exe'
def save_chunks_embeddings(db: Session, item: ChunkEmbeddingCreate, doc_id:int = None):
    text_encode = item.content if item.content else item.title
    text_embedding = embedding_model.encode(text_encode).tolist() # Convert numpy array to list for JSON serialization
    db_item = Document(
        title=item.title,
        content=item.content,
        user_id=item.user_id,
        file_path=item.file_path,
        file_size=item.file_size
    ) if not doc_id else ChunkDocument(
        content=item.content,
        embedding=text_embedding,
        document_id=doc_id
    )
    db.add(db_item)
    db.commit()
    db.refresh(db_item)
    return db_item

def process_csv_file(filepath: str):
    encodings = ['utf-8', 'utf-8-sig', 'windows-1252', 'latin-1']

    df = None
    for encoding in encodings:
        try:
            df = pd.read_csv(filepath, encoding=encoding)
            break
        except (UnicodeDecodeError, TypeError):
            continue

    if df is None:
        raise ValueError("Unable to decode CSV file with specified encodings.")
    
    return df

def extract_text_from_image(filepath: str) -> str:
    try:
        image = Image.open(filepath)
        text = pytesseract.image_to_string(image, lang = "vie+eng")
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code = 500, detail = f"Fail to read image file: {str(e)}")
        
def extract_text_and_chunk(filepath: str, chunk_size: int = 500, chunk_overlap: int = 50):
    chunks = []
    file_extension = filepath.split(".")[-1].lower()
    try:
        if file_extension == "pdf":
            with pdfplumber.open(filepath) as reader:
                for page_nums, page in enumerate(reader.pages, start=1):
                    text = page.extract_text() or ""

                    tables = page.extract_tables()
                    table_texts = []

                    for table in tables:
                        if not table:
                            continue

                        df_table = pd.DataFrame(table)
                        df_table = df_table.fillna("")

                        md_table = df_table.to_markdown(index=False, tablefmt = "github")
                        table_texts.append(f"Data table: \n{md_table}\n")

                    full_page_content = text
                    if table_texts:
                        full_page_content += "".join(table_texts)

                    if not full_page_content.strip():
                        continue

                    start = 0 
                    while start < len(text):
                        end = min(start + chunk_size, len(text))
                        chunk_text = full_page_content[start:end].strip()
                        chunks.append({
                            "text": chunk_text,
                            "page": page_nums
                        })

                        start += (chunk_size - chunk_overlap)

        elif file_extension == "md":
            with open(filepath, "r", encoding="utf-8") as f:
                text = f.read()

            headers_to_split = [
                ('#', 'Header1'),
                ('##', 'Header2'),
                ('###', 'Header3'),
            ]

            # Intialize the splitter with the headers
            markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split=headers_to_split)
            md_header_splitter = markdown_splitter.split_text(text)

            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size = chunk_size,
                chunk_overlap = chunk_overlap,
                seperators = ["\n\n", "\n", " ", ""]
            )

            splits = text_splitter.split_documents(md_header_splitter) # Split the text with headers and content (include page_content and metadata)
            for idx, split in enumerate(splits, start=1):
                chunk_text = split.page_content.strip()

                if not chunk_text:
                    continue
                
                # Create header and context from metadata and split to each other by " > "
                header_content = " > ".join([f"{header}: {content}" for header, content in split.metadata.items()])
                full_chunk_text = chunk_text
                if header_content:
                    full_chunk_text = f"[{header_content}]\n{chunk_text}" #Make content with full context and headers

                chunks.append({
                    "text": full_chunk_text,
                    "source_location": f"{filepath}:chunk:{idx}"
                })
        elif file_extension in ["png", "jpg", "jpeg"]:
            full_image_text = extract_text_from_image(filepath)
            
            start = 0
            idx = 1
            while start < len(full_image_text):
                end = min(start + chunk_size, len(full_image_text))
                chunk_text = full_image_text[start:end].strip()
                
                if chunk_text:
                    chunks.append({
                        "text": chunk_text,
                        "source_location": f"{filepath}:chunk:{idx}"
                    })
                    idx += 1
                
                start += (chunk_size - chunk_overlap)
        elif file_extension in ["txt", "html"]:
            with open(filepath, "r", encoding="utf-8") as f:
                text = f.read()
            start = 0
            while start < len(text):
                end = min(start + chunk_size, len(text))
                chunk_text = text[start:end].strip()
                chunks.append({
                    "text": chunk_text, 
                    "source_location": f"{filepath}:{start}:{end}"
                }) if chunk_text else None
                start += (chunk_size - chunk_overlap)

        elif file_extension == "docx":
            doc = DocxReader(filepath)
            full_text = "\n".join([paragraph.text for paragraph in doc.paragraphs])

            start = 0
            while start < len(full_text):
                end = min(start + chunk_size, len(full_text))
                chunk_text = full_text[start:end].strip()
                chunks.append({
                    "text": chunk_text,
                    "source_location": f"{filepath}:{start}:{end}"
                }) if chunk_text else None
                start += (chunk_size - chunk_overlap)
        
        elif file_extension == "pptx":
            prs = Presentation(filepath)
            for idx, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes: 
                    slide_text.append(shape.text) if hasattr(shape, "text") and shape.text.strip() else None
                text = "\n".join(slide_text)
                if text.strip():
                    chunks.append({
                        "text": text.strip(),
                        "slide": idx
                    })
        
        elif file_extension in ["csv", "xls", "xlsx"]:
            df = process_csv_file(filepath) if file_extension == "csv" else pd.read_excel(filepath)

            chunk_text = []
            rows_per_chunk_text = 25
            for idx, row in df.iterrows():
                row_dict = row.to_dict()
                row_text = ", ".join([f"{col}: {val}" for col, val in row_dict.items() if pd.notna(val)])
                if row_text:
                    chunk_text.append(f"[Row {idx + 1}] {row_text}")

                if (idx + 1) % rows_per_chunk_text == 0 or idx + 1 == len(df):
                    if chunk_text:
                        chunks.append({
                            "text": "\n".join(chunk_text),
                            "row_range": f"{idx + 1 - len(chunk_text)} - {idx + 1}"
                        })
                        chunk_text = []
    except Exception as e:
        raise HTTPException(status_code = 500, detail = f"Fail to read PDF file: {str(e)}")
    
    return chunks


def extract_full_text(filepath: str):
    file_extension = filepath.split(".")[-1].lower()
    full_text = ""
    try:
        if file_extension == "pdf": 
            with pdfplumber.open(filepath) as reader:
                text_list = []
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        text_list.append(text)

                    table_list = page.extract_tables()
                    if table_list:
                        for table in table_list:
                            df_table = pd.DataFrame(table)
                            text_list.append(f"Data table: \n{df_table.to_markdown(index=False)}\n")
                full_text = "\n".join(text_list)
        elif file_extension in ["png", "jpg", "jpeg"]:
            full_text = extract_text_from_image(filepath)
        elif file_extension ==  "docx":
            doc = DocxReader(filepath)
            full_text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        elif file_extension in ["html", "txt", "md"]:
            with open(filepath, "r", encoding = "utf-8") as f: 
                full_text = f.read()
        elif file_extension == "pptx":
            prs = Presentation(filepath)
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    slide_text.append(shape.text) if hasattr(shape, "text") and shape.text.strip() else None
                full_text += "\n".join(slide_text)
        elif file_extension in ["csv", "xls", "xlsx"]:
            df =  process_csv_file(filepath) if file_extension == "csv" else pd.read_excel(filepath)
            
            row_list = []
            for idx, row in df.iterrows():
                row_dict = row.to_dict()
                row_text = ", ".join([f"{col}: {val}" for col, val in row_dict.items() if pd.notna(val)])
                row_list.append(row_text) if row_text else None
            full_text = "\n".join(row_list)
    except Exception as e:
        raise HTTPException(status_code = 500, detail = f"Fail to read file: {str(e)}")
    
    return full_text.strip()

def process_chunks_task(db_factory, document_id: int, file_path: str): # Run in background
    db: Session = db_factory()
    try:
        chunks = extract_text_and_chunk(file_path)
        for chunk in chunks:
            if not chunk['text'].strip():
                continue
            chunk_data = ChunkEmbeddingCreate(
                content=chunk['text'],
                document_id=document_id
            )
            save_chunks_embeddings(db, chunk_data, doc_id=document_id)
            print(f"Background tasks completed for all chunks of document that has id {document_id}")
    except Exception as e:
        print(f"Error during background tasks during processing chunks of document {document_id}: {str(e)}")
    finally:
        db.close()

# Upload file
UPLOAD_DIR = "storage"

def process_conflict_same_title(background_tasks, db: Session, document: Document, user_id: int, filename: str, file_path: str, conflict_strategy: str):
    if conflict_strategy == "skip":
        if os.path.exists(file_path):
            os.remove(file_path) # remove file (already temporarily stored) from storage
        return {
            "id": document.id,
            "title": document.title,
            "file_size": document.file_size,
            "created_at": document.created_at,
            "is_shared": False
        }
    else:
        file_content = extract_full_text(file_path)

        if document.file_path and os.path.exists(document.file_path):
            try:
                os.remove(document.file_path)
            except Exception:
                pass
        
        # Delete all chunks of the document
        db.query(ChunkDocument).filter(ChunkDocument.document_id == document.id).delete()

        # Update the document details
        document.content = file_content
        document.file_path = file_path
        document.file_size = os.path.getsize(file_path)
        document.created_at = datetime.now(timezone.utc)
        db.commit()
        db.refresh(document)
        background_tasks.add_task(process_chunks_task, SessionLocal, document.id, file_path)

        return {
            "id": document.id,
            "title": document.title,
            "file_size": document.file_size,
            "created_at": document.created_at,
            "is_shared": False
        }

def save_loaded_file(db: Session, file: UploadFile, user_id: int, background_tasks, conflict_strategy: str = "rename"):
    user_dir = os.path.join(UPLOAD_DIR, str(user_id))
    os.makedirs(user_dir, exist_ok = True)

    file_extension = os.path.splitext(file.filename)[1]
    file_name = str(uuid.uuid4()) + file_extension
    
    file_path = os.path.join(user_dir, file_name)
    with open(file_path, "wb") as buffer: #Save file into storage
        copyfileobj(file.file, buffer)

    document = db.query(Document).filter(Document.title == file.filename, Document.user_id == user_id).first()

    final_title = file.filename

    if document:
        if conflict_strategy in ["skip", "overwrite"]:
            return process_conflict_same_title(background_tasks, db, document, user_id, file.filename, file_path, conflict_strategy)
        elif conflict_strategy == "rename":
            final_title = resolved_filename_conflict(db, user_id, file.filename)
        else:
            if os.path.exists(file_path):
                os.remove(file_path) # remove file (already temporarily stored) from storage
            raise HTTPException(status_code = 400, detail = "Document with the same name already exists.")
    
    file_content = extract_full_text(file_path)
    
    saved_documents = save_chunks_embeddings(db, ChunkEmbeddingCreate(
        title=final_title, 
        content=file_content,
        user_id=user_id, 
        file_path=file_path, 
        file_size=os.path.getsize(file_path)
    ))
    print("Saved_documents: ", saved_documents.title)
    document_id = saved_documents.id

    background_tasks.add_task(process_chunks_task, SessionLocal, document_id, file_path)

    print("Saved_documents ID: ", saved_documents.id)
    return {
        "id": saved_documents.id,
        "title": saved_documents.title,
        "file_size": saved_documents.file_size,
        "created_at": saved_documents.created_at,
        "is_shared": False
    }



# Mime type Google Workspace to normal mime type
GOOGLE_MIME_TYPE_MAPPING = {
    "application/vnd.google-apps.document": ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx"),
    "application/vnd.google-apps.spreadsheet": ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "xlsx"),
    "application/vnd.google-apps.presentation": ("application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx")
}

def download_file_from_google_drive (file_id: str, access_token: str, user_id: int):
    creds = Credentials(token = access_token)
    services = build('drive', 'v3', credentials = creds)

    file_metadata = services.files().get(fileId = file_id, fields = "name, mimeType").execute()
    original_filename = file_metadata['name']
    mime_type = file_metadata['mimeType']

    user_dir = os.path.join(UPLOAD_DIR, str(user_id))
    if not user_dir: 
        os.makedirs(user_dir, exist_ok = True)

    # if mime_type is google_mime_type (google docx, google sheet, google pptx), then we need to use export_media
    # else we need to use get_media to download directly 
    if mime_type in GOOGLE_MIME_TYPE_MAPPING:
        mime_type, file_extension = GOOGLE_MIME_TYPE_MAPPING[mime_type]
        request = services.files().export_media(fileId = file_id, mimeType = mime_type)
    else: 
        file_extension = os.path.splitext(original_filename)[1]
        request = services.files().get_media(fileId = file_id)

    fileIO = io.BytesIO()
    downloader = MediaIoBaseDownload(fileIO, request)
    done = False

    while not done:
        status, done = downloader.next_chunk()
        print(f"Downloaded {int(status.progress() * 100)}%.")

    fileIO.seek(0) # reset cursor to the beginning of the file

    file_extension = file_extension if not file_extension or file_extension.startswith(".") else f".{file_extension}"
    file_name = str(uuid.uuid4()) + file_extension

    file_path = os.path.join(user_dir, file_name)
    with open(file_path, "wb") as buffer:
        buffer.write(fileIO.read())
        buffer.flush() # Force OS to push all data from cache to disk
        
    fileIO.close() # Free up memory

    return file_path, original_filename

def save_google_drive_file(db: Session, file_id: str, access_token: str, user_id: int, background_tasks, conflict_strategy: str = "rename"):
    file_path, original_filename = download_file_from_google_drive(file_id, access_token, user_id)

    document = db.query(Document).filter(Document.title == original_filename, Document.user_id == user_id).first()

    final_title = original_filename

    if document:
        if conflict_strategy in ["skip", "overwrite"]:
            return process_conflict_same_title(background_tasks, db, document, user_id, original_filename, file_path, conflict_strategy)
        elif conflict_strategy == "rename":
            final_title = resolved_filename_conflict(db, user_id, original_filename)
        else:
            if os.path.exists(file_path):
                os.remove(file_path) # remove file (already temporarily stored) from storage
            raise HTTPException(status_code = 400, detail = "Document with the same name already exists.")
    
    file_content = extract_full_text(file_path)

    saved_documents = save_chunks_embeddings(db, ChunkEmbeddingCreate(
        title = final_title,
        content = file_content,
        file_path = file_path,
        file_size = os.path.getsize(file_path),
        user_id = user_id
    ))

    background_tasks.add_task(process_chunks_task, SessionLocal, saved_documents.id, file_path)

    return {
        "id": saved_documents.id,
        "title": original_filename,
        "file_size": saved_documents.file_size,
        "created_at": saved_documents.created_at,
        "is_shared": False
    }



def get_user_document(db: Session, user_id: int) -> list[Document]:
    return db.query(Document).filter(Document.user_id == user_id, Document.deleted_at == None).all()

async def get_document_insights(db: Session, document_id: int) -> DocumentInsight:
    doc = db.query(Document).filter(Document.id == document_id, Document.deleted_at == None).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found in your account's storage.")
    try: 
        existing_insight = db.query(DocumentInsight).filter(DocumentInsight.document_id == document_id).first()
        if existing_insight and len(existing_insight.key_words) > 0:
            return existing_insight
        insights = await rag_service.generate_document_summary(doc.content)
        saved_insights = DocumentInsight(
            document_id = document_id,
            summary = insights["summary"],
            key_points = list(insights["key_points"]),
            key_words = list(insights["key_words"])
        )
        db.add(saved_insights)
        db.commit()
        return saved_insights
    except Exception as e:
        print(f"Error during getting insights of document {document_id}: {str(e)}")
        return DocumentInsight(
            document_id = document_id,
            summary = "Not summary available",
            key_points = [],
            key_words = []
        )

def delete_document(db: Session, document_ids: list[int], user_id: int):
    docs = db.query(Document).filter(Document.id.in_(document_ids), Document.user_id == user_id, Document.deleted_at == None)
    if not docs:
        raise HTTPException(status_code=404, detail="Document not found in your account's storage.")
    
    docs.update({"deleted_at": datetime.now(timezone.utc)}, synchronize_session = False)
    db.commit()
    return {"message": "Move document to trash successfully!!!"}

def get_trash_document(db: Session, user_id: int):
    return db.query(Document).outerjoin(DocumentShare, Document.id == DocumentShare.document_id)\
        .filter(
            Document.deleted_at != None,
            or_(
                Document.user_id == user_id,
                DocumentShare.shared_to_id == user_id
            )
        )\
        .distinct().all()

def permanent_delete_document(db: Session, document_ids: list[int], user_id: int):
    docs = db.query(Document).filter(Document.id.in_(document_ids), Document.user_id == user_id, Document.deleted_at != None).all()

    if not docs:
        raise HTTPException(status_code=404, detail="Document not found in your account's trash.")
    
    for doc in docs:
        if os.path.exists(doc.file_path):
            try:
                os.remove(doc.file_path) # remove file (already temporarily stored) from storage
            except Exception as e:
                print(f"Error during permanent delete document {doc.id}: {str(e)}")

        db.query(ChunkDocument).filter(ChunkDocument.document_id == doc.id).delete()
        db.query(DocumentShare).filter(DocumentShare.document_id == doc.id).delete()
        db.query(DocumentInsight).filter(DocumentInsight.document_id == doc.id).delete()

        db.delete(doc)
    
    db.commit()
    return {"message": "Delete document and relevant data has been deleted permanently!!!"}

def restore_document(db: Session, document_ids: list[int], user_id: int):
    docs = db.query(Document).filter(Document.id.in_(document_ids), Document.user_id == user_id, Document.deleted_at != None)
    if not docs:
        raise HTTPException(status_code=404, detail="Document not found in your account's trash.")
    
    docs.update({"deleted_at": None}, synchronize_session = False)
    db.commit()
    return {"message": "Restore document successfully!!!"}

async def shared_document_to_user(shared_data: DocumentShareCreate, db: Session, user_id: int):
    sender = db.query(User).filter(User.id == user_id).first()
    if not sender:
        raise HTTPException(status_code=404, detail="Sender not found.")

    doc = db.query(Document).filter(
        Document.id == shared_data.document_id,
        Document.user_id == user_id,
        Document.deleted_at == None
    ).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found in your account's storage.")
    
    receiver = db.query(User).filter(User.email == shared_data.shared_to_email).first()
    if not receiver:
        raise HTTPException(status_code=404, detail="Receiver not found.")
    if receiver.id == user_id:
        raise HTTPException(status_code=400, detail="You cannot share document to yourself.")
    
    existing_share = db.query(DocumentShare).filter(
        DocumentShare.document_id == shared_data.document_id,
        DocumentShare.shared_to_id == receiver.id,
        DocumentShare.shared_by_id == user_id
    ).first()

    if existing_share:
        existing_share.permission = shared_data.permission
        existing_share.created_at = datetime.now(timezone.utc)
        db.commit()
        db.refresh(existing_share)
        share = existing_share
    else: 
        new_share = DocumentShare(
            shared_by_id = user_id,
            shared_to_id = receiver.id,
            document_id = shared_data.document_id,
            permission = shared_data.permission
        )
        doc.shares.append(new_share)
        db.add(new_share)
        db.commit()
        db.refresh(new_share)
        share = new_share

    await trigger_new_notification(db, sender, receiver, doc.title)
    
    return share

async def trigger_new_notification(db: Session, sender: User, receiver: User, doc_title: str):
    notif_receive = Notification(
        user_id = receiver.id,
        text = f"{sender.email} shared a document '{doc_title}' with you.",
        type = "received"
    )

    notif_sent = Notification(
        user_id = sender.id,
        text = f"You shared a document '{doc_title}' with {receiver.email}.",
        type = "sent"
    )

    db.add(notif_receive)
    db.add(notif_sent)
    db.commit()
    db.refresh(notif_receive)
    db.refresh(notif_sent)

    # Emit a signal to send a notification in the background
    await notification_manager.send_notification(
            receiver.id,{
                "id": notif_receive.id,
                "text": notif_receive.text,
                "type": notif_receive.type,
                "created_at": notif_receive.created_at.isoformat()
            }
        )
    
    await notification_manager.send_notification(
            sender.id,{
                "id": notif_receive.id,
                "text": notif_receive.text,
                "type": notif_receive.type,
                "created_at": notif_receive.created_at.isoformat()
            }
        )



def get_shared_document(db: Session, user_id: int, filter_condition):
    document_share = db.query(Document, DocumentShare).join(DocumentShare, Document.id == DocumentShare.document_id).filter(
        filter_condition,
        Document.deleted_at == None
    ).all()
    raw_dict = [{
        "id": doc.id,
        "title": doc.title,
        "shared_by": share.shared_by_id,
        "share_to": share.shared_to_id,
        "permission": share.permission,
        "shared_at": share.created_at.strftime("%Y-%m-%d %H:%M:%S")
    } for doc, share in document_share]

    list_email_share_by = [db.query(User).filter(User.id == dict["shared_by"]).first().email for dict in raw_dict]
    list_email_share_to = [db.query(User).filter(User.id == dict["share_to"]).first().email for dict in raw_dict]

    return [{
        "id": dict["id"],
        "title": dict["title"],
        "shared_by": email_share_by,
        "share_to": email_share_to,
        "permission": dict["permission"],
        "shared_at": dict["shared_at"]
    } for dict, email_share_by, email_share_to in zip(raw_dict, list_email_share_by, list_email_share_to)]
    
def resolved_filename_conflict(db: Session, user_id: int, original_title: str):
    base, extension = os.path.splitext(original_title)
    counter = 1
    new_title = original_title

    while db.query(Document).filter(Document.title == new_title, Document.user_id == user_id).first():
        new_title = f"{base}_({counter}){extension}"
        counter += 1

    return new_title